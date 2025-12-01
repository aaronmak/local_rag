#!/usr/bin/env python3
"""Script to ingest documents from a directory into the RAG system."""

import argparse
from pathlib import Path
from typing import List

from docx import Document as DocxDocument
from pptx import Presentation

from local_rag import RAGPipeline
from local_rag.pdf_processor import PDFLayoutProcessor


def load_pdf_file(file_path: Path, preserve_layout: bool = True) -> tuple[str, dict]:
    """Load a PDF file and extract text with layout awareness.

    Args:
        file_path: Path to PDF file
        preserve_layout: Whether to preserve layout information (bounding boxes, positions, etc.)

    Returns:
        Tuple of (content, metadata)
    """
    processor = PDFLayoutProcessor(preserve_layout=preserve_layout)
    content, metadata = processor.process_pdf(file_path)
    return content, metadata


def load_text_file(file_path: Path) -> tuple[str, dict]:
    """Load a text file.

    Args:
        file_path: Path to text file

    Returns:
        Tuple of (content, metadata)
    """
    content = file_path.read_text(encoding="utf-8")
    metadata = {
        "source": str(file_path),
        "filename": file_path.name,
        "file_type": "txt",
    }

    return content, metadata


def load_docx_file(file_path: Path) -> tuple[str, dict]:
    """Load a Word document (.docx) file.

    Args:
        file_path: Path to Word document

    Returns:
        Tuple of (content, metadata)
    """
    doc = DocxDocument(file_path)

    # Extract text from paragraphs
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

    # Extract text from tables
    table_texts = []
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                table_texts.append(row_text)

    # Combine all content
    content_parts = paragraphs
    if table_texts:
        content_parts.append("\n--- Tables ---\n")
        content_parts.extend(table_texts)

    content = "\n\n".join(content_parts)

    metadata = {
        "source": str(file_path),
        "filename": file_path.name,
        "file_type": "docx",
        "num_paragraphs": len(paragraphs),
        "num_tables": len(doc.tables),
    }

    return content, metadata


def load_pptx_file(file_path: Path) -> tuple[str, dict]:
    """Load a PowerPoint (.pptx) file.

    Args:
        file_path: Path to PowerPoint file

    Returns:
        Tuple of (content, metadata)
    """
    prs = Presentation(file_path)

    slide_contents = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []

        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())

        if slide_texts:
            slide_content = f"--- Slide {slide_num} ---\n" + "\n\n".join(slide_texts)
            slide_contents.append(slide_content)

    content = "\n\n".join(slide_contents)

    metadata = {
        "source": str(file_path),
        "filename": file_path.name,
        "file_type": "pptx",
        "num_slides": len(prs.slides),
    }

    return content, metadata


def load_documents(directory: Path, preserve_layout: bool = True) -> List[tuple[str, dict]]:
    """Load all supported documents from a directory.

    Args:
        directory: Directory containing documents
        preserve_layout: Whether to preserve layout information for PDFs

    Returns:
        List of tuples (content, metadata)
    """
    documents = []
    supported_extensions = {".txt", ".pdf", ".docx", ".pptx"}

    for file_path in directory.rglob("*"):
        if file_path.suffix.lower() not in supported_extensions:
            continue

        try:
            if file_path.suffix.lower() == ".pdf":
                content, metadata = load_pdf_file(file_path, preserve_layout=preserve_layout)
            elif file_path.suffix.lower() == ".docx":
                content, metadata = load_docx_file(file_path)
            elif file_path.suffix.lower() == ".pptx":
                content, metadata = load_pptx_file(file_path)
            else:
                content, metadata = load_text_file(file_path)

            documents.append((content, metadata))
            print(f"Loaded: {file_path}")
        except Exception as e:
            print(f"Error loading {file_path}: {e}")

    return documents


def main():
    """Ingest documents into the RAG system."""
    parser = argparse.ArgumentParser(description="Ingest documents into the RAG system")
    parser.add_argument(
        "directory",
        type=Path,
        help="Directory containing documents to ingest",
    )
    parser.add_argument(
        "--reset",
        action="store_true",
        help="Reset the vector store before ingesting",
    )
    parser.add_argument(
        "--no-layout",
        action="store_true",
        help="Disable layout preservation for PDFs (faster but loses spatial context)",
    )

    args = parser.parse_args()

    if not args.directory.exists():
        print(f"Error: Directory {args.directory} does not exist")
        return

    # Initialize RAG pipeline
    print("Initializing RAG pipeline...")
    rag = RAGPipeline()

    # Reset if requested
    if args.reset:
        print("Resetting vector store...")
        rag.reset()

    # Load documents
    preserve_layout = not args.no_layout
    if preserve_layout:
        print(f"\nLoading documents from {args.directory} (with layout preservation)...")
    else:
        print(f"\nLoading documents from {args.directory} (basic mode)...")

    docs_with_metadata = load_documents(args.directory, preserve_layout=preserve_layout)

    if not docs_with_metadata:
        print("No documents found!")
        return

    # Separate documents and metadata
    documents = [doc for doc, _ in docs_with_metadata]
    metadatas = [meta for _, meta in docs_with_metadata]

    # Add to RAG system
    print(f"\nAdding {len(documents)} documents to the RAG system...")
    ids = rag.add_documents(documents, metadatas)

    # Print stats
    stats = rag.get_stats()
    print("\nIngestion complete!")
    print(f"  Total documents in system: {stats['num_documents']}")
    print(f"  Documents added: {len(ids)}")


if __name__ == "__main__":
    main()
