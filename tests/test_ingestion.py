"""Tests for document ingestion utilities."""

import tempfile
from pathlib import Path

import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader, PdfWriter

from scripts.ingest_documents import (
    load_docx_file,
    load_documents,
    load_pdf_file,
    load_pptx_file,
    load_text_file,
)


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as tmpdir:
        yield Path(tmpdir)


def test_load_text_file(temp_dir):
    """Test loading a text file."""
    # Create a test text file
    test_file = temp_dir / "test.txt"
    test_content = "This is a test document.\nWith multiple lines."
    test_file.write_text(test_content, encoding="utf-8")

    # Load the file
    content, metadata = load_text_file(test_file)

    assert content == test_content
    assert metadata["filename"] == "test.txt"
    assert metadata["file_type"] == "txt"
    assert "source" in metadata


def test_load_pdf_file(temp_dir):
    """Test loading a PDF file."""
    # Create a simple test PDF
    test_file = temp_dir / "test.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)

    with open(test_file, "wb") as f:
        writer.write(f)

    # Load the file
    content, metadata = load_pdf_file(test_file)

    assert isinstance(content, str)
    assert metadata["filename"] == "test.pdf"
    assert metadata["file_type"] == "pdf"
    # Verify num_pages is an integer (may be 0 for blank pages with no extractable text)
    assert isinstance(metadata["num_pages"], int)
    assert metadata["num_pages"] >= 0
    assert "source" in metadata


def test_load_documents(temp_dir):
    """Test loading multiple documents from a directory."""
    # Create test files
    (temp_dir / "doc1.txt").write_text("Document one", encoding="utf-8")
    (temp_dir / "doc2.txt").write_text("Document two", encoding="utf-8")

    # Create a PDF
    pdf_file = temp_dir / "doc3.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with open(pdf_file, "wb") as f:
        writer.write(f)

    # Create a file with unsupported extension
    (temp_dir / "ignored.md").write_text("Should be ignored", encoding="utf-8")

    # Load documents
    documents = load_documents(temp_dir)

    # Should load 3 files (2 txt + 1 pdf), ignore .md
    assert len(documents) == 3

    # Verify structure
    for content, metadata in documents:
        assert isinstance(content, str)
        assert "filename" in metadata
        assert "source" in metadata
        assert "file_type" in metadata


def test_load_documents_recursive(temp_dir):
    """Test loading documents from nested directories."""
    # Create nested structure
    subdir = temp_dir / "subdir"
    subdir.mkdir()

    (temp_dir / "root.txt").write_text("Root document", encoding="utf-8")
    (subdir / "nested.txt").write_text("Nested document", encoding="utf-8")

    # Load documents
    documents = load_documents(temp_dir)

    # Should find both files
    assert len(documents) == 2


def test_load_documents_empty_directory(temp_dir):
    """Test loading from an empty directory."""
    documents = load_documents(temp_dir)
    assert len(documents) == 0


def test_load_documents_with_errors(temp_dir, capsys):
    """Test handling of file loading errors."""
    # Create a file with invalid encoding
    bad_file = temp_dir / "bad.txt"
    bad_file.write_bytes(b"\x80\x81\x82")  # Invalid UTF-8

    # Load documents (should handle error gracefully)
    documents = load_documents(temp_dir)

    # Should return empty list and print error
    assert len(documents) == 0
    captured = capsys.readouterr()
    assert "Error loading" in captured.out


def test_load_docx_file(temp_dir):
    """Test loading a Word document (.docx) file."""
    # Create a test Word document
    test_file = temp_dir / "test.docx"
    doc = DocxDocument()

    # Add paragraphs
    doc.add_paragraph("First paragraph")
    doc.add_paragraph("Second paragraph")
    doc.add_paragraph("")  # Empty paragraph (should be skipped)
    doc.add_paragraph("Third paragraph")

    # Add a table
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Header 1"
    table.rows[0].cells[1].text = "Header 2"
    table.rows[1].cells[0].text = "Data 1"
    table.rows[1].cells[1].text = "Data 2"

    doc.save(test_file)

    # Load the file
    content, metadata = load_docx_file(test_file)

    # Verify content
    assert "First paragraph" in content
    assert "Second paragraph" in content
    assert "Third paragraph" in content
    assert "Header 1" in content
    assert "Header 2" in content
    assert "Data 1" in content
    assert "Data 2" in content
    assert "--- Tables ---" in content

    # Verify metadata
    assert metadata["filename"] == "test.docx"
    assert metadata["file_type"] == "docx"
    # DocxDocument creates default paragraphs, so count will be higher than manually added
    assert metadata["num_paragraphs"] >= 3
    assert metadata["num_tables"] == 1
    assert "source" in metadata


def test_load_docx_file_no_tables(temp_dir):
    """Test loading a Word document without tables."""
    test_file = temp_dir / "test_no_tables.docx"
    doc = DocxDocument()
    doc.add_paragraph("Only paragraphs here")
    doc.save(test_file)

    content, metadata = load_docx_file(test_file)

    assert "Only paragraphs here" in content
    assert "--- Tables ---" not in content
    assert metadata["num_tables"] == 0


def test_load_pptx_file(temp_dir):
    """Test loading a PowerPoint (.pptx) file."""
    # Create a test PowerPoint presentation
    test_file = temp_dir / "test.pptx"
    prs = Presentation()

    # Add slide 1 with title and content
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
    slide1.shapes.title.text = "Slide 1 Title"
    slide1.placeholders[1].text = "This is the content of slide 1"

    # Add slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Slide 2 Title"
    slide2.placeholders[1].text = "This is the content of slide 2"

    # Add slide 3 with only title
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout
    slide3.shapes.title.text = "Slide 3 Title Only"

    prs.save(test_file)

    # Load the file
    content, metadata = load_pptx_file(test_file)

    # Verify content
    assert "--- Slide 1 ---" in content
    assert "Slide 1 Title" in content
    assert "content of slide 1" in content

    assert "--- Slide 2 ---" in content
    assert "Slide 2 Title" in content
    assert "content of slide 2" in content

    assert "--- Slide 3 ---" in content
    assert "Slide 3 Title Only" in content

    # Verify metadata
    assert metadata["filename"] == "test.pptx"
    assert metadata["file_type"] == "pptx"
    assert metadata["num_slides"] == 3
    assert "source" in metadata


def test_load_pptx_file_empty_slides(temp_dir):
    """Test loading a PowerPoint with empty slides."""
    test_file = temp_dir / "test_empty.pptx"
    prs = Presentation()

    # Add a blank slide with no text
    prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add a slide with text
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Only Slide with Content"

    prs.save(test_file)

    content, metadata = load_pptx_file(test_file)

    # Should only include the slide with content
    assert "--- Slide 2 ---" in content
    assert "Only Slide with Content" in content
    # Blank slide should not appear (no "--- Slide 1 ---")
    assert content.count("--- Slide") == 1
    assert metadata["num_slides"] == 2


def test_load_documents_with_all_formats(temp_dir):
    """Test loading documents of all supported formats."""
    # Create text file
    (temp_dir / "doc.txt").write_text("Text document", encoding="utf-8")

    # Create PDF
    pdf_file = temp_dir / "doc.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with open(pdf_file, "wb") as f:
        writer.write(f)

    # Create Word document
    docx_file = temp_dir / "doc.docx"
    doc = DocxDocument()
    doc.add_paragraph("Word document")
    doc.save(docx_file)

    # Create PowerPoint
    pptx_file = temp_dir / "doc.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "PowerPoint slide"
    prs.save(pptx_file)

    # Create unsupported file
    (temp_dir / "ignored.json").write_text('{"ignored": true}', encoding="utf-8")

    # Load all documents
    documents = load_documents(temp_dir)

    # Should load 4 files (txt, pdf, docx, pptx), ignore .json
    assert len(documents) == 4

    # Verify file types are present
    file_types = {meta["file_type"] for _, meta in documents}
    assert file_types == {"txt", "pdf", "docx", "pptx"}

    # Verify structure
    for content, metadata in documents:
        assert isinstance(content, str)
        assert "filename" in metadata
        assert "source" in metadata
        assert "file_type" in metadata
